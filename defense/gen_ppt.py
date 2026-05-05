#!/usr/bin/env python3
"""Generate thesis defense PPT from ppt-plan.md using python-pptx."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import re

HUST_BLUE = RGBColor(0, 63, 136)
HUST_BLUE_LIGHT = RGBColor(0, 102, 178)
DARK_TEXT = RGBColor(51, 51, 51)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(240, 240, 240)
ACCENT_ORANGE = RGBColor(230, 126, 34)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def set_text(
    shape,
    text,
    font_size=18,
    bold=False,
    color=DARK_TEXT,
    alignment=PP_ALIGN.LEFT,
    font_name="微软雅黑",
):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(
    tf,
    text,
    font_size=18,
    bold=False,
    color=DARK_TEXT,
    alignment=PP_ALIGN.LEFT,
    space_before=Pt(4),
    space_after=Pt(2),
    font_name="微软雅黑",
    level=0,
):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    p.level = level
    return p


def add_title_bar(slide, title_text, subtitle_text=None):
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), fill_color=HUST_BLUE)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11), Inches(0.7))
    set_text(
        txBox, title_text, font_size=32, bold=True, color=WHITE, font_name="微软雅黑"
    )
    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.75), Inches(11), Inches(0.4)
        )
        set_text(
            txBox2,
            subtitle_text,
            font_size=16,
            bold=False,
            color=RGBColor(200, 220, 240),
            font_name="微软雅黑",
        )
    add_shape(
        slide, Inches(0), Inches(1.2), SLIDE_W, Inches(0.06), fill_color=ACCENT_ORANGE
    )


def add_bullet_content(
    slide, bullets, top=Inches(1.6), left=Inches(0.8), width=Inches(11.5), font_size=20
):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(bullet, str):
            text, fs, b, c = bullet, font_size, False, DARK_TEXT
        elif len(bullet) == 4:
            text, fs, b, c = bullet
        elif len(bullet) == 3:
            text, fs, b = bullet
            c = DARK_TEXT
        elif len(bullet) == 2:
            text, fs = bullet
            b = False
            c = DARK_TEXT
        else:
            text = bullet[0]
            fs = font_size
            b = False
            c = DARK_TEXT
        p.text = text
        p.font.size = Pt(fs if fs else font_size)
        p.font.bold = b
        p.font.color.rgb = c if c else DARK_TEXT
        p.font.name = "微软雅黑"
        p.space_before = Pt(8)
        p.space_after = Pt(4)
    return txBox


def add_two_column(slide, left_bullets, right_bullets, top=Inches(1.6)):
    col_w = Inches(5.5)
    add_bullet_content(slide, left_bullets, top=top, left=Inches(0.8), width=col_w)
    add_bullet_content(slide, right_bullets, top=top, left=Inches(7), width=col_w)


def add_note(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# ====== Slide 1: Cover ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=HUST_BLUE)
add_shape(
    slide, Inches(0), Inches(3.2), SLIDE_W, Inches(0.08), fill_color=ACCENT_ORANGE
)

txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.0), Inches(10), Inches(1.0))
set_text(
    txBox,
    "本科毕业论文答辩",
    font_size=24,
    bold=False,
    color=RGBColor(180, 200, 230),
    alignment=PP_ALIGN.CENTER,
)

txBox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.333), Inches(1.4))
set_text(
    txBox,
    "基于驱动程序替换的物联网设备\n固件仿真技术研究",
    font_size=36,
    bold=True,
    color=WHITE,
    alignment=PP_ALIGN.CENTER,
)

txBox = slide.shapes.add_textbox(Inches(2), Inches(3.8), Inches(9), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ""
lines = [
    "答辩人：XXX          指导教师：XXX  教授",
    "",
    "华中科技大学  XXX  学院",
    "",
    "2026 年 X 月 X 日",
]
for line in lines:
    p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.name = "微软雅黑"
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(6)
add_note(slide, "（不讲，翻页）")

# ====== Slide 2: Outline ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "汇报提纲")
outline_items = [
    ("01", "研究背景与问题"),
    ("02", "方法概述与主要贡献"),
    ("03", "系统设计（三个核心模块）"),
    ("04", "系统实现"),
    ("05", "实验验证"),
    ("06", "总结与展望"),
]
y_start = Inches(1.8)
for i, (num, title) in enumerate(outline_items):
    y = y_start + Inches(i * 0.8)
    box = add_shape(
        slide, Inches(1.5), y, Inches(0.8), Inches(0.6), fill_color=HUST_BLUE
    )
    set_text(box, num, font_size=24, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    box.text_frame.auto_size = None
    txBox = slide.shapes.add_textbox(
        Inches(2.6), y + Inches(0.08), Inches(8), Inches(0.5)
    )
    set_text(txBox, title, font_size=24, bold=False, color=DARK_TEXT)
add_note(
    slide,
    "各位老师好，我是XXX，我的论文题目是《基于驱动程序替换的物联网设备固件仿真技术研究》，接下来从六个部分进行汇报。",
)

# ====== Slide 3: Background ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "研究背景与核心问题")
add_bullet_content(
    slide,
    [
        ("研究背景", 22, True, HUST_BLUE),
        "  物联网设备数量爆发式增长，固件安全问题日益突出",
        "  （Mirai 僵尸网络 DDoS 攻击、俄乌冲突 Industroyer 电力系统攻击等）",
        "  固件仿真（Firmware Re-Hosting）是进行固件安全测试的有效手段",
        "",
        ("核心问题：自动化程度不足，人工成本高昂", 22, True, ACCENT_ORANGE),
        "  基于硬件建模：需为每种外设手动建立仿真模型，工作量大且通用性差",
        "  基于驱动层替换（如 HALucinator）：函数定位、语义理解、替换代码生成仍高度依赖人工干预",
        "  缺乏自动化反馈机制：仿真异常时无法自动诊断和修复，依赖分析人员经验",
    ],
    font_size=19,
)
add_note(
    slide,
    "首先是研究背景。随着物联网设备数量爆发式增长，固件安全问题日益突出，近年来 Mirai 僵尸网络和俄乌冲突中的电力系统攻击都表明固件安全形势严峻。固件仿真是进行安全测试的有效手段，但现有方法的自动化程度严重不足。具体来说，基于硬件建模的方案需要为每种外设手动建立仿真模型，工作量大且通用性差；基于驱动层替换的方案虽然避免了底层建模，但函数的定位、语义理解和替换代码生成仍然高度依赖人工干预。此外，现有方法大多缺乏自动化的反馈机制，仿真过程中出现异常时无法自动诊断和修复。因此，本文致力于解决固件仿真技术自动化程度不足这一核心问题。",
)

# ====== Slide 4: Comparison ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "现有固件仿真技术对比")
# Table
rows, cols = 4, 4
table_shape = slide.shapes.add_table(
    rows, cols, Inches(0.8), Inches(1.6), Inches(11.5), Inches(3.5)
)
table = table_shape.table
headers = ["方案", "典型工具", "优点", "缺点"]
data = [
    ["基于硬件建模", "Pretender, sEmu", "模拟精度高", "需要真实设备，通用性差"],
    [
        "基于程序分析",
        "P2IM, Fuzzware",
        "自动推断部分行为",
        "推断规则有限，难以覆盖复杂外设",
    ],
    [
        "基于驱动层替换",
        "HALucinator,\nSAFIREFUZZ",
        "降低硬件依赖，速度快",
        "函数定位与替换依赖人工，\n自动化不足",
    ],
]
for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "微软雅黑"
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = HUST_BLUE
for i, row_data in enumerate(data):
    for j, val in enumerate(row_data):
        cell = table.cell(i + 1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.name = "微软雅黑"
            if j == 0:
                p.font.bold = True
        if i == 2:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(230, 240, 255)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(5.4), Inches(11.5), Inches(0.6))
set_text(
    txBox,
    "驱动函数的定位与替换是所有方案的共同瓶颈",
    font_size=20,
    bold=True,
    color=ACCENT_ORANGE,
    alignment=PP_ALIGN.CENTER,
)
add_note(
    slide,
    "我们将现有方案分为三类进行对比。基于硬件建模的方案需要真实设备观测，通用性差；基于程序分析的方案能自动推断部分行为，但规则有限。基于驱动层替换的方案在降低硬件依赖和提升速度方面具有优势，但驱动函数的定位与替换仍然高度依赖人工。可以看到，无论采用哪种技术路线，驱动函数的自动化定位与替换都是核心瓶颈，这也是本文重点解决的问题。",
)

# ====== Slide 5: Method Overview ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "本文方法概述", "三环节可衔接的自动化替换流程")
modules = [
    ("静态分析", "自动识别硬件依赖函数\n（MMIO 模式识别 + CodeQL 查询）"),
    ("LLM 驱动替换", "函数语义分类 → 按类别生成替换代码\n→ 编译修复循环"),
    ("动态反馈闭环", "仿真运行 → 异常诊断\n→ 增量修正 → 回归验证"),
]
for i, (title, desc) in enumerate(modules):
    x = Inches(1.0 + i * 4.0)
    box = add_shape(
        slide,
        x,
        Inches(1.8),
        Inches(3.2),
        Inches(2.8),
        fill_color=HUST_BLUE if i == 1 else RGBColor(70, 130, 180),
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "微软雅黑"
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(12)
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(220, 230, 240)
    p2.font.name = "微软雅黑"
    p2.alignment = PP_ALIGN.CENTER
    if i < 2:
        arrow = slide.shapes.add_textbox(
            x + Inches(3.2), Inches(2.8), Inches(0.8), Inches(0.6)
        )
        set_text(
            arrow,
            "→",
            font_size=36,
            bold=True,
            color=ACCENT_ORANGE,
            alignment=PP_ALIGN.CENTER,
        )

txBox = slide.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(11), Inches(0.8))
set_text(
    txBox,
    "核心目标：在可控迭代中把固件从\u201c可构建\u201d推进到\u201c可稳定运行\u201d",
    font_size=20,
    bold=True,
    color=ACCENT_ORANGE,
    alignment=PP_ALIGN.CENTER,
)
txBox2 = slide.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(11), Inches(0.5))
set_text(
    txBox2,
    "配图：论文图3.1 系统整体架构图",
    font_size=14,
    bold=False,
    color=RGBColor(150, 150, 150),
    alignment=PP_ALIGN.CENTER,
)
add_note(
    slide,
    "针对上述问题，本文提出一种基于驱动程序替换的自动化固件仿真方法。系统将整个流程拆分为三个可衔接环节。第一个环节是静态分析，通过 MMIO 模式识别和 CodeQL 查询自动定位硬件依赖函数；第二个环节是基于 LLM 的驱动替换，对识别出的函数进行语义分类，按类别生成替换代码，并通过编译修复循环保障代码可编译通过；第三个环节是动态反馈闭环，在仿真运行中发现问题时自动诊断、修正并回归验证。核心目标不是一次性得到完美替换，而是在可控迭代中把固件从\u201c可构建\u201d逐步推进到\u201c可稳定运行\u201d。整体架构如右图所示。",
)

# ====== Slide 6: Contributions ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "主要贡献")
contributions = [
    (
        "(1) 硬件依赖代码自动识别",
        [
            "提出三种 MMIO 定位模式（直接地址强转、间接初始化、基址直传）",
            "基于 CodeQL 查询自动定位驱动函数，构建可查询驱动知识库",
            "覆盖 STM32、NXP、Atmel 三类平台 + FreeRTOS、RT-Thread、Zephyr、裸机四种系统",
        ],
    ),
    (
        "(2) 基于 LLM 的驱动函数语义理解与替换",
        [
            "提出七类函数分类体系，设计四种替换策略（Skip / ReturnOK / Redirect / Event）",
            "通过结构化提示词模板引导 LLM 生成语义保持的替代代码",
        ],
    ),
    (
        "(3) 动态测试反馈闭环",
        [
            "设计\u201c运行\u2192诊断\u2192修复\u2192回归\u201d自动化闭环，三个智能体协同调度",
            "集成 AFL++ 模糊测试，复用仿真环境与注入接口",
        ],
    ),
]
y = Inches(1.5)
for title, items in contributions:
    txBox = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.5), Inches(0.4))
    set_text(txBox, title, font_size=20, bold=True, color=HUST_BLUE)
    y += Inches(0.45)
    for item in items:
        txBox = slide.shapes.add_textbox(Inches(1.2), y, Inches(11), Inches(0.35))
        set_text(txBox, "  • " + item, font_size=16, color=DARK_TEXT)
        y += Inches(0.35)
    y += Inches(0.15)
add_note(
    slide,
    "本文的主要贡献有三点。第一，提出了硬件依赖代码的自动识别方法，通过三种 MMIO 定位模式和 CodeQL 查询，覆盖了三种芯片平台和四种操作系统。第二，探索性地将大语言模型应用于驱动函数替换任务，设计了七类分类体系和四种替换策略，通过结构化提示词引导 LLM 生成语义保持的替代代码。第三，构建了动态测试反馈闭环，通过三个智能体协同调度实现自动化的迭代优化，并集成 AFL++ 进行模糊测试验证。",
)

# ====== Slide 7: MMIO Pattern Recognition ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "模块一：MMIO 访问模式识别", "驱动函数接口静态分析与识别")
add_bullet_content(
    slide,
    [
        ("嵌入式固件通过内存映射 I/O（MMIO）与外设交互", 20, True, HUST_BLUE),
        "  寄存器映射到统一地址空间，通过普通内存读写指令访问",
        "  核心问题：如何从源码中自动定位所有 MMIO 访问点？",
        "",
        ("三种 MMIO 定位模式", 22, True, ACCENT_ORANGE),
        ("1. 直接地址强转", 18, True, HUST_BLUE),
        "  GPIOA->ODR = val  —  基址直接硬编码在代码中",
        ("2. 间接初始化", 18, True, HUST_BLUE),
        "  ptr = &USART1; ptr->BRR = val  —  基址通过变量间接传递",
        ("3. 基址直传", 18, True, HUST_BLUE),
        "  HAL_UART_Init(&huart1)  —  基址作为函数参数传入",
    ],
    font_size=17,
)
add_note(
    slide,
    "首先介绍第一个模块——静态分析与识别。嵌入式固件通过 MMIO 与外设交互，外设寄存器被映射到统一地址空间。核心问题是如何从源码中自动定位所有 MMIO 访问点。我们提出三种定位模式来覆盖不同的基址传递方式：直接地址强转是最常见的模式，基址直接硬编码在代码中；间接初始化模式中，基址通过全局变量或函数参数间接传递；基址直传模式中，外设基址直接作为函数参数传入。通过这三种模式，系统能够覆盖固件中绝大多数的 MMIO 访问方式。",
)

# ====== Slide 8: CodeQL ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "模块一：CodeQL 查询与数据模型")
add_two_column(
    slide,
    [
        ("CodeQL 驱动函数定位", 22, True, HUST_BLUE),
        "  基于三种 MMIO 模式编写 CodeQL 查询规则",
        "  自动提取：",
        "    MMIO 结构体实例",
        "    → 寄存器字段访问",
        "    → 函数调用链",
        "",
        "  数据流追踪：",
        "  将外设基址常量作为源点",
        "  沿赋值和参数传播追踪到字段解引用",
    ],
    [
        ("统一数据模型", 22, True, HUST_BLUE),
        ("  FunctionInfo", 16, True),
        "    函数签名、调用关系、MMIO 访问集合",
        ("  MMIOAccess", 16, True),
        "    寄存器地址、访问类型（读/写/读-改-写）",
        ("  ProjectAnalysisResult", 16, True),
        "    工程级驱动函数汇总",
        "",
        ("驱动函数定位准则", 18, True, ACCENT_ORANGE),
        "  函数内部直接包含 MMIO 访问，",
        "  或调用含 MMIO 访问的函数",
    ],
)
add_note(
    slide,
    "基于这三种模式，我们编写了 CodeQL 查询规则来定位驱动函数。查询逻辑分为三步：首先通过结构体类型名定位 MMIO 实例，然后追踪寄存器字段级别的读写访问，最后沿函数调用链和数据流传播确定哪些函数涉及硬件依赖。所有分析结果被组织为统一的数据模型，包括函数信息、MMIO 访问记录和工程级汇总结果，供后续 LLM 模块消费。右图展示了分析智能体的工作流程。",
)

# ====== Slide 9: Classification ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "模块二：七类函数分类体系", "核心范式：先分类后替换")
rows, cols = 8, 4
table_shape = slide.shapes.add_table(
    rows, cols, Inches(0.6), Inches(1.5), Inches(12), Inches(5.0)
)
table = table_shape.table
headers = ["类别", "语义", "替换方式", "典型函数"]
data = [
    ["INIT", "外设初始化配置", "Skip（移除寄存器写入）", "HAL_UART_Init"],
    ["RECV", "从硬件读取数据", "Redirect（仿真注入数据）", "HAL_UART_Receive"],
    ["LOOP", "轮询等待硬件状态", "Skip（移除死循环等待）", "HAL_UART_WaitOnFlag"],
    ["IRQ", "中断处理/回调", "Event（改为事件回调）", "UART_IRQHandler"],
    ["PUREDRV", "纯驱动操作", "Skip / ReturnOK", "HAL_GPIO_DeInit"],
    ["CORE", "OS 内核函数", "仿真器直接支持（不替换）", "SysTick_Handler"],
    ["NODRIVER", "纯软件逻辑", "不替换", "main"],
]
for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "微软雅黑"
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = HUST_BLUE
for i, row_data in enumerate(data):
    for j, val in enumerate(row_data):
        cell = table.cell(i + 1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.name = "微软雅黑"
            if j == 0:
                p.font.bold = True
                p.font.color.rgb = HUST_BLUE
            p.alignment = PP_ALIGN.CENTER
        if i % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(245, 248, 255)
add_note(
    slide,
    "第二个模块的核心思想是先分类后替换。我们将驱动函数分为七类。INIT 类是外设初始化配置函数，只需移除寄存器写入即可，采用 Skip 策略；RECV 类负责从硬件读取数据，需要通过仿真接口注入模拟数据，采用 Redirect 策略；LOOP 类包含轮询等待硬件状态的循环，直接移除即可避免仿真死循环；IRQ 类是中断处理函数，改造为可由仿真器主动触发的事件回调；PUREDRV 类是不影响上层逻辑的纯驱动操作，根据返回类型选择 Skip 或 ReturnOK；CORE 类是操作系统内核函数，由仿真器直接支持；NODRIVER 类是纯软件逻辑函数，无需替换。从饼图可以看到，INIT 类占 39%，是硬件依赖最集中的部分；PUREDRV 占 24.3%；NODRIVER 占 20.2%，这些函数经过分析确认不需要硬件解耦。",
)

# ====== Slide 10: Strategies ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "模块二：替换策略与验证机制")
rows, cols = 5, 4
table_shape = slide.shapes.add_table(
    rows, cols, Inches(0.6), Inches(1.5), Inches(12), Inches(2.8)
)
table = table_shape.table
headers = ["策略", "适用类别", "核心思路", "仿真接口"]
data = [
    ["Skip", "INIT, LOOP", "移除函数体，保留调用路径", "无需注入"],
    ["ReturnOK", "PUREDRV", "保留签名，返回成功状态值", "无需注入"],
    [
        "Redirect",
        "RECV",
        "硬件操作替换为 HAL_BE_* 注入",
        "HAL_BE_In,\nHAL_BE_ENET_ReadFrame",
    ],
    ["Event", "IRQ", "改造为可触发的事件回调", "HAL_BE_* + 事件触发"],
]
for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "微软雅黑"
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = HUST_BLUE
for i, row_data in enumerate(data):
    for j, val in enumerate(row_data):
        cell = table.cell(i + 1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.name = "微软雅黑"
            if j == 0:
                p.font.bold = True
                p.font.color.rgb = HUST_BLUE
            p.alignment = PP_ALIGN.CENTER

add_bullet_content(
    slide,
    [
        ("语义保持性三条件", 20, True, ACCENT_ORANGE),
        ("  1. 调用正确性", 17, True),
        "     替换函数签名与原函数一致，调用方无需修改",
        ("  2. 返回正确性", 17, True),
        "     所有执行路径的返回值与声明类型兼容",
        ("  3. 可观测行为一致性", 17, True),
        "     固件到达关键执行节点，无崩溃/死循环/断言失败",
    ],
    top=Inches(4.6),
    font_size=16,
)
add_note(
    slide,
    "针对需要替换的五类函数，我们设计了四种替换策略。重点介绍 Redirect 策略：它将原始的硬件操作替换为对 HAL_BE_* 注入接口的调用，仿真器在运行时通过这些接口向固件注入模拟数据。例如，HAL_UART_Receive 被替换后，当固件调用该函数时，实际执行的是 HAL_BE_In 接口，仿真器将测试数据写入固件的接收缓冲区。为保障替换质量，我们定义了语义保持性的三个递进条件：调用正确性确保调用方无需修改，返回正确性确保返回值类型兼容，可观测行为一致性确保固件能到达关键执行节点。",
)

# ====== Slide 11: Prompt Engineering ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "模块二：LLM 提示词工程")
add_two_column(
    slide,
    [
        ("分类提示词（四层结构）", 22, True, HUST_BLUE),
        ("  1. 角色定义", 17, True),
        "     嵌入式软件工程师",
        "     消除驱动函数的硬件依赖",
        ("  2. 工具说明", 17, True),
        "     GetFunctionInfo",
        "     → GetMMIOFunctionInfo",
        "     → [按需调用其他工具]",
        ("  3. 分类策略", 17, True),
        "     七类函数的识别特征",
        "     替换策略和关键约束",
        ("  4. 输出格式", 17, True),
        "     结构化 JSON",
        "     （类别、理由、替换代码）",
    ],
    [
        ("修复提示词", 22, True, HUST_BLUE),
        "",
        ("  先取证再修复", 17, True),
        "  要求 LLM 先调用诊断工具",
        "  获取完整执行流和错误信息",
        "",
        ("  上下文完整获取", 17, True),
        "  修复前必须获取该函数的",
        "  历次替换历史",
        "",
        ("  步数预算控制", 17, True),
        "  限制最多分析 3-5 个",
        "  可疑函数后立即生成修复方案",
    ],
)
add_note(
    slide,
    "提示词工程是 LLM 模块的关键技术。分类提示词采用四层结构：首先将 LLM 定位为嵌入式软件工程师，明确任务是消除硬件依赖；然后声明五个可调用的工具函数及其调用顺序，要求 LLM 先获取函数源码再确认硬件访问方式；接着嵌入七类函数的识别特征和替换策略作为分类依据；最后要求输出结构化 JSON 格式。修复提示词强调先取证再修复的原则，要求 LLM 先获取完整的执行流和错误信息，并且获取该函数的历次替换历史以避免重复失败方案。此外，通过限制步数预算避免 LLM 无效调用工具。",
)

# ====== Slide 12: Feedback Loop ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "模块三：动态测试反馈闭环", "工作流：运行 → 诊断 → 修复 → 回归")
add_two_column(
    slide,
    [
        ("三个智能体分工", 22, True, HUST_BLUE),
        "",
        ("  仿真执行智能体（主控）", 17, True),
        "  加载固件、执行仿真、判断运行状态",
        "  触发诊断/修复/回归",
        "",
        ("  闭环修复智能体", 17, True),
        "  基于错误日志分析问题函数",
        "  生成增量修复方案",
        "",
        ("  闭环构建智能体", 17, True),
        "  应用修复代码，执行编译构建",
        "  返回编译状态",
    ],
    [
        ("终止条件", 22, True, HUST_BLUE),
        "",
        ("  成功终止", 17, True, RGBColor(39, 174, 96)),
        "  固件到达预设关键节点且无崩溃",
        "",
        ("  失败终止", 17, True, RGBColor(231, 76, 60)),
        "  连续 3 轮无新修复（停滞检测）",
        "  达到 10 轮上限",
        "  不可修复错误类型",
        "",
        "",
        ("  配图：论文动态测试反馈闭环工作流图", 14, False, RGBColor(150, 150, 150)),
    ],
)
add_note(
    slide,
    "第三个模块是动态测试反馈闭环，也是本文的重要创新点之一。工作流如左图所示：仿真执行智能体作为主控，加载固件并执行仿真。如果仿真成功，流程结束；如果仿真失败，智能体自动分析日志定位问题函数，触发闭环修复智能体生成修复方案；修复后由闭环构建智能体重新编译构建，再次进行回归验证。整个过程自动循环，直到固件稳定运行。终止条件包括三种情况：成功到达关键节点、连续 3 轮无新修复即判定为停滞、或达到 10 轮迭代上限。",
)

# ====== Slide 13: Simulation Platform ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "模块三：仿真平台与注入接口")
add_two_column(
    slide,
    [
        ("仿真平台", 22, True, HUST_BLUE),
        "  基于 Unicorn ARM 模拟器",
        "  轻量级、可编程钩子机制",
        "  支持 ARMv7-M Thumb 指令集",
        "",
        "  弱符号占位函数实现注入接口",
        "  运行时由宿主机 Python 处理器接管",
    ],
    [
        ("7 种 HAL_BE_* 注入接口", 22, True, HUST_BLUE),
        "  HAL_BE_return_0：恒定返回 0",
        "  HAL_BE_In：定长接收（UART）",
        "  HAL_BE_ENET_ReadFrame：变长帧读取（ETH）",
        "  HAL_BE_Block_Read/Write：块设备读写",
        "",
        ("  同一套接口同时用于：", 17, True, ACCENT_ORANGE),
        "  替换正确性动态验证",
        "  AFL++ 模糊测试输入注入",
    ],
)
add_note(
    slide,
    "仿真平台基于 Unicorn 引擎构建，选择 Unicorn 而非 QEMU 的原因是它在裸机 Cortex-M 场景下更轻量且可控。注入接口通过弱符号占位函数实现：链接阶段占位函数提供默认定义，运行时由宿主机 Python 处理器接管实际执行。表中列出了 7 种注入接口，覆盖了 UART 定长接收、以太网变长帧读取、块设备读写等主要外设场景。需要特别说明的是，同一套注入接口同时用于两方面的验证：一是替换正确性的动态验证，二是模糊测试的输入注入，无需为模糊测试额外搭建仿真环境。",
)

# ====== Slide 14: Implementation ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "系统实现")
add_two_column(
    slide,
    [
        ("技术栈", 22, True, HUST_BLUE),
        ("  静态分析", 17, True),
        "  CodeQL（MMIO 查询 + 数据流追踪）",
        ("  智能体框架", 17, True),
        "  LangGraph（状态机工作流 + 检查点机制）",
        ("  仿真引擎", 17, True),
        "  Unicorn（ARMv7-M Thumb 指令集模拟）",
        ("  模糊测试", 17, True),
        "  AFL++（通过 UnicornAFL 扩展集成）",
    ],
    [
        ("关键实现", 22, True, HUST_BLUE),
        "",
        "  CodeQL 查询规则：",
        "  覆盖直接地址强转、间接初始化、基址直传",
        "",
        "  LangGraph 状态机：",
        "  分析→分类→生成→验证→编译修复",
        "",
        "  UnicornAFL 集成：",
        "  fork-server 协议 + 共享内存位图",
        "  完成覆盖率反馈",
    ],
)
add_note(
    slide,
    "实现层面，系统采用四个核心技术栈。CodeQL 负责静态分析，我们为三种 MMIO 模式分别编写了查询规则；LangGraph 用于构建智能体的图式工作流，通过状态机和条件边实现分析、分类、生成、验证的自动化流程；Unicorn 作为仿真引擎，通过钩子机制实现函数拦截和数据注入；AFL++ 通过 UnicornAFL 扩展与仿真环境集成，实现覆盖率引导的模糊测试。",
)

# ====== Slide 15: Experiment Design ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "实验设计")
add_two_column(
    slide,
    [
        ("测试覆盖范围", 22, True, HUST_BLUE),
        ("  芯片平台", 17, True),
        "  STM32、NXP、Atmel",
        ("  操作系统", 17, True),
        "  FreeRTOS、RT-Thread、Zephyr、裸机",
        ("  测试用例", 17, True),
        "  35 个固件（官方 SDK + 开源社区）",
        ("  外设类型", 17, True),
        "  UART、I2C、以太网、Flash 文件系统",
    ],
    [
        ("评估维度", 22, True, HUST_BLUE),
        "",
        "  1. 驱动函数识别与分类统计",
        "  2. 替换功能正确性",
        "     （编译/启动成功率、运行稳定性）",
        "  3. 动态反馈闭环有效性",
        "  4. 方法效率对比",
        "  5. 模糊测试验证",
        "     （基本块覆盖率对比）",
    ],
)
add_note(
    slide,
    "实验设计覆盖了 STM32、NXP、Atmel 三类芯片平台和 FreeRTOS、RT-Thread、Zephyr、裸机四种操作系统，共 35 个测试固件。固件来源包括芯片厂商官方 SDK 示例和开源社区项目。我们从五个维度进行评估：驱动函数识别与分类、替换功能正确性、反馈闭环有效性、方法效率对比，以及模糊测试验证。",
)

# ====== Slide 16: Recognition Stats ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "实验结果：驱动函数识别与分类统计")
add_two_column(
    slide,
    [
        ("识别统计", 22, True, HUST_BLUE),
        "",
        "  35 个测试用例共识别 4,029 个驱动函数",
        "  （平均 115.1 个/固件）",
        "",
        ("  1,580 个（39.2%）", 18, True, ACCENT_ORANGE),
        "  生成了替换代码",
        "  （平均 45.1 个/固件）",
        "",
        ("  517 个替换代码", 18, True, ACCENT_ORANGE),
        "  经过动态反馈闭环迭代优化",
        "  （平均 14.8 个/固件）",
    ],
    [
        ("分类分布", 22, True, HUST_BLUE),
        "",
        "  INIT:  39.0%  （初始化配置）",
        "  PUREDRV:  24.3%  （纯驱动操作）",
        "  NODRIVER:  20.2%  （纯软件逻辑）",
        "  CORE:  5.1%  （OS 内核）",
        "  IRQ:  ~4%  （中断处理）",
        "  LOOP:  ~4%  （轮询等待）",
        "",
        "  INIT 类硬件依赖最集中",
        "  NODRIVER 类无需替换",
    ],
)
add_note(
    slide,
    "35 个测试用例共识别出 4029 个驱动函数，其中 1580 个生成了替换代码，占比约 39.2%。从分类分布来看，INIT 类占 39%，是硬件依赖最集中的部分；PUREDRV 占 24.3%，这类函数只需简单的 Skip 或 ReturnOK 处理；NODRIVER 占 20.2%，说明有相当比例的函数经过分析确认不需要硬件解耦。还有 517 个替换代码经过了动态反馈闭环的迭代优化，平均每个固件约 15 个。",
)

# ====== Slide 17: Functional Correctness ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "实验结果：替换功能正确性")
add_bullet_content(
    slide,
    [
        ("验证方法", 22, True, HUST_BLUE),
        "  执行轨迹验证：检查关键节点函数是否按预期顺序执行",
        "  行为验证：对比替换前后的 API 调用序列",
        "",
        ("关键指标（14 个代表性样本）", 22, True, ACCENT_ORANGE),
    ],
    font_size=18,
)

metrics = [
    ("编译成功率", "95.3%", "首次生成 79.2%，经编译修复循环提升"),
    ("启动成功率", "100%", "所有固件均成功启动"),
    ("运行稳定性", "无崩溃", "所有固件在仿真环境中稳定运行"),
    ("关键节点匹配率", "100%", "系统启动、外设初始化、数据收发、任务调度"),
]
y = Inches(4.0)
rows_t, cols_t = 5, 3
table_shape = slide.shapes.add_table(
    rows_t, cols_t, Inches(0.8), y, Inches(11.5), Inches(2.8)
)
table = table_shape.table
for j, h in enumerate(["指标", "结果", "说明"]):
    cell = table.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "微软雅黑"
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = HUST_BLUE
for i, (metric, result, note) in enumerate(metrics):
    for j, val in enumerate([metric, result, note]):
        cell = table.cell(i + 1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.name = "微软雅黑"
            if j == 1:
                p.font.bold = True
                p.font.color.rgb = ACCENT_ORANGE
            p.alignment = PP_ALIGN.CENTER
add_note(
    slide,
    "功能正确性从两个维度验证。一是执行轨迹验证，检查系统启动、外设初始化、数据收发等关键节点是否按预期顺序执行；二是行为验证，对比替换前后的 API 调用序列。14 个代表性样本的关键节点匹配率达到 100%。编译成功率 95.3%，其中首次生成即可编译的占 79.2%，剩余部分通过编译修复循环迭代修正。启动成功率和运行稳定性均达到 100%。",
)

# ====== Slide 18: Feedback Loop Effectiveness ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "实验结果：动态反馈闭环有效性")
add_bullet_content(
    slide,
    [
        (
            "32 个用例（91.4%）触发了至少一次替换更新，累计 517 次",
            19,
            True,
            ACCENT_ORANGE,
        ),
        "",
        ("155 个函数属于\u201c闭环新增覆盖\u201d", 20, True, HUST_BLUE),
        "  静态分析阶段分类为 NODRIVER 或未被识别",
        "  仿真运行中暴露出硬件依赖问题",
        "  由反馈闭环补充生成替换代码",
        "  若去除动态反馈闭环，这 155 个函数将缺失替换，对应执行路径受阻",
        "",
        ("典型案例：NXP_I2C_BareMetal", 20, True, HUST_BLUE),
        "  CLOCK_SetDiv 被静态分析分类为 NODRIVER（纯业务逻辑）",
        "  实际包含 while 循环轮询硬件握手标志位 → 仿真中死循环",
        "  反馈闭环自动识别并修复，移除死循环等待逻辑",
    ],
    font_size=17,
)
add_note(
    slide,
    "反馈闭环的有效性体现在两个数据上：91.4% 的测试用例触发了至少一次替换更新，累计 517 次。特别值得关注的是，有 155 个函数属于闭环新增覆盖。这些函数在静态分析阶段被分类为不需要替换或未被识别，但在仿真运行中暴露出硬件依赖问题。例如 NXP_I2C_BareMetal 中的 CLOCK_SetDiv 函数，静态分析将其分类为纯业务逻辑函数，但实际包含轮询硬件握手标志位的 while 循环，在仿真中必然死循环。反馈闭环自动检测到这一问题并生成修复。这说明动态反馈对于弥补静态分析的不足起到了关键作用。",
)

# ====== Slide 19: Efficiency Comparison ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "实验结果：方法效率对比")
rows, cols = 6, 4
table_shape = slide.shapes.add_table(
    rows, cols, Inches(0.8), Inches(1.5), Inches(11.5), Inches(3.5)
)
table = table_shape.table
headers = ["指标", "本文方法", "HALucinator", "Para-rehosting"]
data = [
    ["平均处理驱动函数数", "96.4", "14.6", "16.6"],
    ["平均人工介入时间", "0.6 h", "3.0 h", "4.9 h"],
    ["人工介入降低", "-80%", "-", "-"],
    ["驱动函数数量提升", "5.8× / 6.6×", "-", "-"],
    ["编译成功率", "95.3%", "100%", "100%"],
]
for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "微软雅黑"
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = HUST_BLUE
for i, row_data in enumerate(data):
    for j, val in enumerate(row_data):
        cell = table.cell(i + 1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(15)
            p.font.name = "微软雅黑"
            if j == 1:
                p.font.bold = True
                p.font.color.rgb = HUST_BLUE
            p.alignment = PP_ALIGN.CENTER
        if i % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(245, 248, 255)

add_bullet_content(
    slide,
    [
        "LLM 自动处理时长 24.6h（8 个固件），不占用研究人员工作时间",
        "LLM 处理可自动化、支持并行批处理，且不占用研究人员持续在线时间",
    ],
    top=Inches(5.3),
    font_size=16,
)
add_note(
    slide,
    "与方法效率对比，可以看到三个关键指标。首先，本文方法平均处理 96.4 个驱动函数，是 HALucinator 的 6.6 倍和 Para-rehosting 的 5.8 倍。其次，人工介入时间仅为 0.6 小时每固件，降低约 80%。编译成功率 95.3% 略低于手工方法的 100%，但考虑到显著的自动化优势完全可以接受。需要说明的是，LLM 自动处理时长为 24.6 小时，这部分时间由 LLM 自动执行，不占用研究人员的实际工作时间。",
)

# ====== Slide 20: Fuzzing ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "实验结果：模糊测试验证", "AFL++ × 9 个固件 × 24 小时")
rows, cols = 4, 3
table_shape = slide.shapes.add_table(
    rows, cols, Inches(1.5), Inches(1.6), Inches(10), Inches(2.5)
)
table = table_shape.table
headers = ["指标", "本文方法", "HALucinator"]
data = [
    ["平均基本块覆盖数量", "652.3", "579.2"],
    ["平均执行速度", "29.4 exec/s", "29.5 exec/s"],
    ["基本块增量", "+73.1 (+12.6%)", "-"],
]
for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "微软雅黑"
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = HUST_BLUE
for i, row_data in enumerate(data):
    for j, val in enumerate(row_data):
        cell = table.cell(i + 1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(16)
            p.font.name = "微软雅黑"
            if j == 1 and i == 2:
                p.font.bold = True
                p.font.color.rgb = ACCENT_ORANGE
            p.alignment = PP_ALIGN.CENTER

add_bullet_content(
    slide,
    [
        "9 个样本中 7 个覆盖率正增长，2 个轻微负增长",
        "ETH 类固件优势突出（平均 +209 个基本块），UART/I2C 类基本持平",
        "网络协议栈代码量大、路径复杂，全面替换使模糊测试能探索更多协议处理路径",
    ],
    top=Inches(4.5),
    font_size=17,
)
add_note(
    slide,
    "为验证替换后固件在实际安全测试场景中的可用性，我们选取 9 个覆盖 ETH、UART、I2C 三种外设的固件，使用 AFL++ 进行 24 小时模糊测试。结果表明，本文方法的平均基本块覆盖数量为 652.3 个，较 Halucinator 的 579.2 个增加 73.1 个，提升 12.6%。ETH 类固件的覆盖率提升最为显著，平均增加约 209 个基本块，因为网络协议栈代码量大、路径复杂，本文方法的全面替换使模糊测试能探索更多协议处理路径。UART 和 I2C 类固件基本持平，因为这些外设代码量较小，关键路径已被覆盖。",
)

# ====== Slide 21: Summary ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "总结与展望")
add_two_column(
    slide,
    [
        ("三个创新点", 22, True, HUST_BLUE),
        "",
        ("  1. 硬件依赖代码自动识别", 17, True),
        "  三种 MMIO 模式 + CodeQL 查询",
        "",
        ("  2. LLM 驱动函数语义理解与替换", 17, True),
        "  七类分类 + 四种策略",
        "",
        ("  3. 动态测试反馈闭环", 17, True),
        "  自动化迭代 + 模糊测试集成",
    ],
    [
        ("局限性与未来工作", 22, True, HUST_BLUE),
        "",
        ("  局限性", 17, True, RGBColor(231, 76, 60)),
        "  依赖源码",
        "  平台覆盖有限",
        "  复杂驱动函数精度有待提升",
        "",
        ("  未来工作", 17, True, RGBColor(39, 174, 96)),
        "  支持无源码固件",
        "  扩展平台覆盖",
        "  提升复杂驱动精度",
    ],
)
add_note(
    slide,
    "总结一下。本文针对固件仿真自动化程度不足的问题，提出了三个方面的创新：一是通过三种 MMIO 模式和 CodeQL 查询实现硬件依赖代码的自动识别；二是利用 LLM 进行函数语义理解与替换代码生成，设计了七类分类体系和四种替换策略；三是构建了动态反馈闭环实现自动化的迭代优化。当前方法的局限性在于依赖源码、平台覆盖有限。未来将围绕无源码固件支持、平台扩展和复杂驱动精度提升三个方向继续推进。",
)

# ====== Slide 22: Thanks ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=HUST_BLUE)
add_shape(
    slide, Inches(0), Inches(3.4), SLIDE_W, Inches(0.06), fill_color=ACCENT_ORANGE
)

txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1.2))
set_text(
    txBox, "致  谢", font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER
)

txBox = slide.shapes.add_textbox(Inches(2), Inches(3.8), Inches(9), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ""
for line in [
    "感谢导师 XXX 教授的悉心指导与耐心帮助",
    "",
    "感谢各位评审专家的宝贵意见",
    "",
    "感谢家人和同学朋友的支持",
]:
    p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.name = "微软雅黑"
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(8)
add_note(slide, "以上是我的全部汇报，感谢各位老师的评审，请各位老师批评指正。")

# ====== Backup A: Classification Rules ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "备用A：七类函数分类判定规则", "分类判定规则细节")
rows, cols = 8, 2
table_shape = slide.shapes.add_table(
    rows, cols, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5)
)
table = table_shape.table
headers = ["类别", "判定条件"]
data = [
    ["INIT", "函数名含 Init/Config/Setup；多个 MMIO 寄存器顺序写；位于启动阶段"],
    ["RECV", "函数名含 Receive/Read/Get/Input；含 MMIO 数据寄存器读；参数含缓冲区指针"],
    [
        "LOOP",
        "函数名含 Wait/Poll/Loop；以 MMIO 状态寄存器为条件的 while 循环；循环体无实质数据处理",
    ],
    [
        "IRQ",
        "函数名含 IRQHandler/Callback/Handler；注册为中断服务例程；含中断状态清标志",
    ],
    [
        "PUREDRV",
        "仅含简单 MMIO 操作；返回状态枚举或 void；函数名含 DeInit/MspDeInit/Flush",
    ],
    ["CORE", "OS 内核函数（rt_thread_init, SysTick_Handler）；由仿真器直接提供"],
    ["NODRIVER", "无直接 MMIO 访问；纯软件逻辑（参数整理、句柄映射、状态转发）"],
]
for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "微软雅黑"
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = HUST_BLUE
for i, row_data in enumerate(data):
    for j, val in enumerate(row_data):
        cell = table.cell(i + 1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.name = "微软雅黑"
            if j == 0:
                p.font.bold = True
                p.font.color.rgb = HUST_BLUE
                p.alignment = PP_ALIGN.CENTER
        if i % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(245, 248, 255)
add_note(
    slide,
    "如果被问到分类的具体判定依据，可以展示此页。判定规则从四个维度综合判断：函数名关键词、MMIO 访问模式、控制流结构和调用上下文。",
)

# ====== Backup B: Redirect Example ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "备用B：Redirect 策略代码示例", "HAL_UART_Receive")
add_two_column(
    slide,
    [
        ("原始函数", 20, True, RGBColor(231, 76, 60)),
        "",
        "HAL_StatusTypeDef HAL_UART_receive(",
        "    UART_HandleTypeDef *huart,",
        "    uint8_t *pData,",
        "    uint16_t Size,",
        "    uint32_t Timeout) {",
        "  // ... MMIO 寄存器操作 ...",
        "  while (__HAL_UART_GET_FLAG(",
        "      huart, UART_FLAG_RXNE)",
        "      == RESET);",
        "  *pData = (uint8_t)(",
        "      huart->Instance->RDR & 0xFF);",
        "  return HAL_OK;",
        "}",
    ],
    [
        ("替换函数", 20, True, RGBColor(39, 174, 96)),
        "",
        "HAL_StatusTypeDef HAL_UART_receive(",
        "    UART_HandleTypeDef *huart,",
        "    uint8_t *pData,",
        "    uint16_t Size,",
        "    uint32_t Timeout) {",
        "",
        "  HAL_BE_In(pData, Size);",
        "",
        "  return HAL_OK;",
        "}",
        "",
        "",
        "  仿真器通过 HAL_BE_In",
        "  将测试数据写入固件缓冲区",
    ],
)
add_note(
    slide,
    "这里展示一个 Redirect 策略的具体代码示例。原始的 HAL_UART_Receive 函数通过轮询 RXNE 标志位等待硬件数据，在仿真环境中标志位永远不会置位，会导致死循环。替换后调用 HAL_BE_In 注入接口，由仿真器将测试数据写入固件缓冲区，函数正常返回。",
)

# ====== Backup C: Feedback Case ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(
    slide, "备用C：反馈闭环典型案例", "NXP_I2C_BareMetal — 静态误分类的运行时修正"
)
add_bullet_content(
    slide,
    [
        ("问题描述", 22, True, HUST_BLUE),
        "  CLOCK_SetDiv 被静态分析分类为 NODRIVER（纯业务逻辑）",
        "  实际包含 while 循环轮询 CCM->CDHIPR 硬件握手标志位",
        "  仿真中标志位永远不置位 → 死循环",
        "",
        ("修复过程", 22, True, HUST_BLUE),
        "  反馈闭环自动诊断 → 识别 LOOP 模式 → 移除 while 循环 → 修复成功",
        "",
        ("启示", 22, True, ACCENT_ORANGE),
        "  静态分析阶段的 LLM 分类存在误判率",
        "  反馈闭环作为运行时保障机制有效弥补这一不足",
    ],
    font_size=18,
)
add_note(
    slide,
    "这里举一个典型案例说明反馈闭环的价值。NXP_I2C_BareMetal 中的 CLOCK_SetDiv 函数，静态分析将其分类为不需要替换的纯业务逻辑函数。但实际代码中包含一个 while 循环轮询硬件握手标志位，在仿真环境中该标志位不会被置位，必然死循环。反馈闭环在仿真运行中检测到这一问题，LLM 分析后识别出 LOOP 模式特征，将 while 循环体注释掉，同时保留寄存器写入操作以保持时钟配置语义完整。这说明静态分析存在误判率，反馈闭环作为运行时保障机制是必要的。",
)

# ====== Backup D: LLM Cost ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "备用D：LLM 计算资源消耗")
rows, cols = 5, 2
table_shape = slide.shapes.add_table(
    rows, cols, Inches(2), Inches(1.6), Inches(9), Inches(3.0)
)
table = table_shape.table
headers = ["指标", "数值"]
data = [
    ["8 个固件 LLM 总处理时长", "24.6 小时（平均 3.1 小时/固件）"],
    ["总词元消耗", "约 1.04 亿"],
    ["平均每函数 LLM 处理时长", "约 116 秒"],
    ["平均每函数词元消耗", "约 13.5 万"],
]
for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "微软雅黑"
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = HUST_BLUE
for i, row_data in enumerate(data):
    for j, val in enumerate(row_data):
        cell = table.cell(i + 1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(15)
            p.font.name = "微软雅黑"
            p.alignment = PP_ALIGN.CENTER
        if i % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(245, 248, 255)

add_bullet_content(
    slide,
    [
        "词元消耗与驱动函数数量正相关",
        "LLM 处理为自动化执行时间，不占用研究人员工作时间",
        "未来可通过模型量化、缓存机制等进一步降低成本",
    ],
    top=Inches(5.0),
    font_size=17,
)
add_note(
    slide,
    "如果被问到 LLM 的使用成本，可以展示此页。8 个固件的 LLM 总处理时长为 24.6 小时，总词元消耗约 1.04 亿，平均每个函数需要 116 秒。词元消耗与固件的驱动函数数量正相关。需要注意的是，这部分时间由 LLM 自动执行，不占用研究人员的实际工作时间，且可以通过模型量化、缓存机制等进一步优化成本。",
)

output_path = "/Users/jie/Documents/workspace/my_thesis/defense/thesis_defense.pptx"
prs.save(output_path)
print(f"PPT saved to {output_path}")
